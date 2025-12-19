from flask import Flask, request, render_template, jsonify
import os
from werkzeug.utils import secure_filename
import json
import requests

# Document processing libraries
from docx import Document
from pptx import Presentation
import PyPDF2

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'}

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs('templates', exist_ok=True)

# Model information with advantages
MODEL_INFO = {
    'llama3.2:3b': {
        'name': 'Llama 3.2 3B',
        'size': '2 GB',
        'speed': 'Fast',
        'advantages': 'Best overall balance - Fast and accurate JSON generation',
        'recommended': True,
        'emoji': '🏆'
    },
    'llama3.2:1b': {
        'name': 'Llama 3.2 1B',
        'size': '700 MB',
        'speed': 'Very Fast',
        'advantages': 'Lightning fast processing - Perfect for quick mindmaps',
        'recommended': False,
        'emoji': '⚡'
    },
    'qwen2.5:3b': {
        'name': 'Qwen 2.5 3B',
        'size': '2 GB',
        'speed': 'Fast',
        'advantages': 'Highest accuracy - Excellent at structured data extraction',
        'recommended': False,
        'emoji': '🎯'
    },
    'gemma2:2b': {
        'name': 'Gemma 2 2B',
        'size': '1.6 GB',
        'speed': 'Fast',
        'advantages': 'Lightweight and efficient - Good for educational content',
        'recommended': False,
        'emoji': '💎'
    },
    'phi3:mini': {
        'name': 'Phi-3 Mini',
        'size': '2.3 GB',
        'speed': 'Fast',
        'advantages': 'Compact powerhouse - Microsoft\'s efficient model',
        'recommended': False,
        'emoji': '🚀'
    },
    'mistral:7b': {
        'name': 'Mistral 7B',
        'size': '4.1 GB',
        'speed': 'Moderate',
        'advantages': 'Premium quality - Best for detailed complex mindmaps',
        'recommended': False,
        'emoji': '👑'
    }
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_docx(file_path):
    """Extract text from DOCX files - comprehensive extraction"""
    doc = Document(file_path)
    text = []
    
    # Extract from paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text.strip())
    
    # Extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                if cell.text.strip():
                    row_data.append(cell.text.strip())
            if row_data:
                text.append(" | ".join(row_data))
    
    return '\n'.join(text)

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint files - ALL content from ALL slides"""
    prs = Presentation(file_path)
    text = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text.append(f"SLIDE {slide_num}:")
        
        # Extract text from all shapes
        for shape in slide.shapes:
            # Text in shapes
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
            
            # Text in text frames
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text.append(paragraph.text.strip())
            
            # Text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_data.append(cell.text.strip())
                    if row_data:
                        text.append(" | ".join(row_data))
        
        text.append("")  # Add spacing between slides
    
    return '\n'.join(text)

def extract_text_from_pdf(file_path):
    """Extract text from PDF files"""
    text = []
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text.append(f"PAGE {page_num}:")
                    text.append(page_text)
                    text.append("")
    except Exception as e:
        print(f"PDF extraction error: {e}")
    
    return '\n'.join(text)

def extract_text(file_path):
    """Extract text based on file type"""
    ext = file_path.rsplit('.', 1)[1].lower()
    
    try:
        if ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext in ['doc', 'docx']:
            return extract_text_from_docx(file_path)
        elif ext in ['ppt', 'pptx']:
            return extract_text_from_pptx(file_path)
        elif ext == 'pdf':
            return extract_text_from_pdf(file_path)
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""
    
    return ""

def generate_mindmap_with_ollama(content, model="llama3.2:3b"):
    """Generate comprehensive mindmap structure using Ollama API"""
    
    # Keep more content for better context (but not too much)
    content_preview = content[:5000] if len(content) > 5000 else content
    
    prompt = f"""You are a mindmap generation assistant. Analyze the following document/presentation content and create a COMPREHENSIVE hierarchical mindmap structure.

CONTENT TO ANALYZE:
{content_preview}

TASK:
Create a detailed JSON mindmap with multiple levels of hierarchy. 
- Extract the MAIN TOPIC from the content
- Identify 4-7 MAJOR SUBTOPICS (not just 3-5)
- Add 3-5 KEY POINTS under each subtopic (not just 2-4)
- Make the structure RICH and DETAILED

REQUIRED JSON FORMAT:
{{
    "name": "Main Topic Name",
    "children": [
        {{
            "name": "Subtopic 1",
            "children": [
                {{"name": "Key Point 1"}},
                {{"name": "Key Point 2"}},
                {{"name": "Key Point 3"}},
                {{"name": "Key Point 4"}}
            ]
        }},
        {{
            "name": "Subtopic 2",
            "children": [
                {{"name": "Key Point 1"}},
                {{"name": "Key Point 2"}},
                {{"name": "Key Point 3"}}
            ]
        }}
    ]
}}

CRITICAL RULES:
1. Keep names concise but descriptive (max 60 characters)
2. Extract ACTUAL topic names and concepts from content
3. Create 4-7 main subtopics (not just 2-3)
4. Add 3-5 detail points per subtopic (not just 2)
5. Organize hierarchically and logically
6. Return ONLY valid JSON, no markdown, no explanations
7. Ensure proper JSON formatting with quotes and commas
8. Do not use trailing commas
9. Make the mindmap RICH with details from all slides/pages

JSON OUTPUT:"""

    try:
        print(f"🤖 Calling Ollama with model: {model}")
        print(f"📊 Content size: {len(content_preview)} characters")
        
        # Call Ollama API using requests
        response = requests.post(
            'http://localhost:11434/api/generate',
            json={
                'model': model,
                'prompt': prompt,
                'stream': False,
                'options': {
                    'temperature': 0.3,
                    'top_p': 0.9,
                    'num_predict': 2000  # Allow longer responses
                }
            },
            timeout=180  # Increased timeout
        )
        
        if response.status_code == 200:
            result = response.json()
            response_text = result.get('response', '').strip()
            print(f"📝 Ollama response length: {len(response_text)} chars")
            
            # Try to extract JSON from response
            response_text = response_text.replace('``````', '').strip()
            
            # Find JSON object
            start_idx = response_text.find('{')
            end_idx = response_text.rfind('}') + 1
            
            if start_idx != -1 and end_idx > start_idx:
                json_str = response_text[start_idx:end_idx]
                
                # Clean up trailing commas
                json_str = json_str.replace(',\n]', '\n]').replace(',\n}', '\n}')
                json_str = json_str.replace(', ]', ']').replace(', }', '}')
                
                try:
                    mindmap_data = json.loads(json_str)
                    
                    # Validate structure
                    if 'name' in mindmap_data and 'children' in mindmap_data:
                        # Count nodes
                        node_count = count_nodes(mindmap_data)
                        print(f"✅ Valid mindmap generated with {node_count} nodes")
                        return mindmap_data
                except json.JSONDecodeError as je:
                    print(f"❌ JSON parsing error: {je}")
                    print(f"Error position: {je.pos}, Line: {je.lineno}, Col: {je.colno}")
            
            raise ValueError("Invalid JSON structure in response")
            
    except requests.exceptions.ConnectionError:
        print("❌ Cannot connect to Ollama")
        return {
            "error": "Cannot connect to Ollama. Make sure Ollama is running.",
            "fallback": True
        }
    except requests.exceptions.Timeout:
        print("❌ Ollama request timeout")
        return {
            "error": "Request timeout. Try a smaller file or faster model.",
            "fallback": True
        }
    except Exception as e:
        print(f"❌ Error calling Ollama: {e}")
    
    # Fallback: Create detailed structure from content
    print("⚠️ Using fallback mindmap (AI generation failed)")
    return create_detailed_fallback_mindmap(content)

def count_nodes(data):
    """Count total nodes in mindmap"""
    count = 1
    if 'children' in data:
        for child in data['children']:
            count += count_nodes(child)
    return count

def create_detailed_fallback_mindmap(content):
    """Create a detailed mindmap when AI fails - extract all content"""
    lines = [line.strip() for line in content.split('\n') if line.strip() and len(line.strip()) > 3]
    
    # Extract main topic
    main_topic = "Document Content"
    for line in lines[:5]:
        if len(line) > 10 and not line.startswith('SLIDE') and not line.startswith('PAGE'):
            main_topic = line[:60]
            break
    
    # Group lines into meaningful chunks
    children = []
    current_section = None
    section_items = []
    
    for line in lines:
        # Skip slide/page markers
        if line.startswith('SLIDE') or line.startswith('PAGE'):
            if current_section and section_items:
                children.append({
                    "name": current_section,
                    "children": [{"name": item[:55]} for item in section_items[:5]]
                })
            current_section = line.replace('SLIDE ', '').replace('PAGE ', '').replace(':', '')
            section_items = []
        elif current_section and len(line) > 5:
            section_items.append(line)
    
    # Add last section
    if current_section and section_items:
        children.append({
            "name": current_section,
            "children": [{"name": item[:55]} for item in section_items[:5]]
        })
    
    # If no children from section parsing, create generic structure
    if not children:
        chunk_size = max(4, len(lines) // 5)
        for i in range(0, min(len(lines), 20), chunk_size):
            chunk = lines[i:i+chunk_size]
            if chunk:
                subtopic_name = chunk[0][:50]
                details = [{"name": line[:50]} for line in chunk[1:6]]
                
                if details:
                    children.append({
                        "name": subtopic_name,
                        "children": details
                    })
    
    # Ensure we have children
    if not children:
        children = [
            {
                "name": "Overview",
                "children": [
                    {"name": "Main Points"},
                    {"name": "Key Concepts"},
                    {"name": "Details"}
                ]
            }
        ]
    
    return {
        "name": main_topic,
        "children": children[:7]  # Limit to 7 main topics
    }

@app.route('/')
def index():
    return render_template('mindmap.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        try:
            # Extract text
            print(f"\n📄 Extracting text from {filename}...")
            content = extract_text(filepath)
            
            if not content.strip():
                if os.path.exists(filepath):
                    os.remove(filepath)
                return jsonify({'error': 'No text content found in file. File may be corrupted or empty.'}), 400
            
            print(f"✅ Extracted {len(content)} characters")
            print(f"📊 Preview: {content[:200]}...")
            
            # Get model from request
            model = request.form.get('model', 'llama3.2:3b')
            
            # Generate mindmap
            mindmap_data = generate_mindmap_with_ollama(content, model)
            
            # Check for errors BEFORE deleting file
            if mindmap_data.get('error'):
                if os.path.exists(filepath):
                    os.remove(filepath)
                return jsonify({
                    'error': mindmap_data['error'],
                    'mindmap': mindmap_data if mindmap_data.get('fallback') else None
                }), 500
            
            # Get model info
            model_info = MODEL_INFO.get(model, {
                'name': model,
                'advantages': 'Custom model'
            })
            
            # Clean up uploaded file AFTER successful generation
            if os.path.exists(filepath):
                os.remove(filepath)
            
            return jsonify({
                'success': True,
                'mindmap': mindmap_data,
                'content_length': len(content),
                'model_used': model,
                'model_name': model_info.get('name', model),
                'model_advantages': model_info.get('advantages', '')
            })
            
        except Exception as e:
            if os.path.exists(filepath):
                os.remove(filepath)
            print(f"❌ Processing error: {e}")
            return jsonify({'error': f'Processing error: {str(e)}'}), 500
    
    return jsonify({'error': 'Invalid file type. Supported: TXT, PDF, DOC, DOCX, PPT, PPTX'}), 400

@app.route('/models', methods=['GET'])
def get_models():
    """Get available Ollama models with their info"""
    try:
        response = requests.get('http://localhost:11434/api/tags', timeout=5)
        if response.status_code == 200:
            installed_models = response.json().get('models', [])
            installed_model_names = [m['name'] for m in installed_models]
            
            # Return model info for installed models
            available_models = []
            for model_name in MODEL_INFO.keys():
                model_data = MODEL_INFO[model_name].copy()
                model_data['id'] = model_name
                model_data['installed'] = model_name in installed_model_names
                available_models.append(model_data)
            
            return jsonify({'models': available_models})
    except Exception as e:
        print(f"Error fetching models: {e}")
    
    # Return default model info if can't connect
    default_models = []
    for model_name, info in MODEL_INFO.items():
        model_data = info.copy()
        model_data['id'] = model_name
        model_data['installed'] = False
        default_models.append(model_data)
    
    return jsonify({'models': default_models})

if __name__ == '__main__':
    print("=" * 60)
    print("🚀 Starting Mindmap Generator...")
    print("📍 Make sure Ollama is running in background")
    print("🌐 Access at: http://localhost:5002")
    print("=" * 60)
    app.run(debug=True, host='0.0.0.0', port=5002)
